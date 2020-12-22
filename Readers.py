import docx2txt
from pptx import Presentation
import PyPDF2
import xlrd
from itertools import chain

class FileReader():

    def __init__(self, filename):

        self.filename = filename
        self.path = '/'.join(filename.split("/")[0:-1])
        self.name = filename.split("/")[-1]
        self.extension = filename.split("/")[-1].split(".")[-1]
        self.text = self.extractText()

    def extractText(self):

        text = None
        if self.extension in ("docx", "doc"):
            text = self.read_doc()
        elif self.extension == "pdf":
            text = self.read_pdf()
        elif self.extension == "txt":
            text = self.read_text()
        elif self.extension == "ppt" or self.extension == "pptx":
            text = self.read_ppt()
        elif self.extension in ("xlsx", "xls"):
            text = self.read_excel()
        else:
            raise Exception("ERROR: Extension not implemented")
        return text


    def read_pdf(self):

        with open(self.filename, "rb") as f:
            reader = PyPDF2.PdfFileReader(f)
            page = reader.getPage(0)
            text = page.extractText()
        return text

    def read_doc(self):

        text = docx2txt.process(self.filename)
        text = text.rstrip()

        return text
    
    def read_excel(self):

        workbook = xlrd.open_workbook(self.filename)

        text_list = []
        cols_name = []
        for s in range(workbook.nsheets):
            ws = workbook.sheet_by_index(s)
            cols_name.append([str(c).upper() for c in ws.row_values(0)])
            for row in ws.get_rows():
                for r in row:
                    text_list.append(str(r.value))
            text = ' '.join(text_list)

        return text,  list(chain.from_iterable(cols_name))

    def read_text(self):

        f = open(self.filename, "r", encoding="latin-1")
        return f.read()


    def read_ppt(self):

        def extract_shape_txt(shape):
            if hasattr(shape, "text"):
                shapetext = shape.text
                return shapetext
            else:
                return ''

        prs = Presentation(self.filename)
        text = ' '
        for slide in prs.slides:
            for shape in slide.shapes:
                if "GroupShape" in shape.__str__():
                    for innershape in shape.shapes:
                        text += extract_shape_txt(innershape)
                else:
                    text += extract_shape_txt(shape)

        return text


    def read_image(self):
        pass

